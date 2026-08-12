from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

# Initialize Widescreen 16:9 Presentation
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Color Palette Definitions
DARK_BG = RGBColor(11, 15, 25)       # Obsidian Dark Slate (#0B0F19)
CYAN_ACCENT = RGBColor(0, 240, 255)   # Neon Cyan (#00F0FF)
EMERALD_ACCENT = RGBColor(16, 185, 129)# Neon Emerald (#10B981)
TEXT_LIGHT = RGBColor(241, 245, 249) # Off-White Text (#F1F5F9)
CARD_BG = RGBColor(22, 30, 46)       # Dark Panel Background (#161E2E)
MUTED_BORDER = RGBColor(30, 41, 59)   # Slate Border (#1E293B)

def add_blank_slide_with_bg(prs):
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)
    bg_shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg_shape.fill.solid()
    bg_shape.fill.fore_color.rgb = DARK_BG
    bg_shape.line.fill.background()
    return slide

def add_header(slide, title_text, category_text="AGENTX SANDBOX"):
    cat_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11.7), Inches(0.4))
    tf_cat = cat_box.text_frame
    tf_cat.word_wrap = True
    p_cat = tf_cat.paragraphs[0]
    p_cat.text = category_text.upper()
    p_cat.font.size = Pt(11)
    p_cat.font.bold = True
    p_cat.font.color.rgb = CYAN_ACCENT
    
    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.7), Inches(11.7), Inches(0.8))
    tf_title = title_box.text_frame
    tf_title.word_wrap = True
    p_title = tf_title.paragraphs[0]
    p_title.text = title_text
    p_title.font.size = Pt(26)
    p_title.font.bold = True
    p_title.font.color.rgb = TEXT_LIGHT

# SLIDE 1: Title Slide
slide1 = add_blank_slide_with_bg(prs)
title_box = slide1.shapes.add_textbox(Inches(1.0), Inches(2.2), Inches(11.3), Inches(3.0))
tf1 = title_box.text_frame
tf1.word_wrap = True

p1 = tf1.paragraphs[0]
p1.text = "AgentX Sandbox"
p1.font.size = Pt(54)
p1.font.bold = True
p1.font.color.rgb = CYAN_ACCENT

p2 = tf1.add_paragraph()
p2.text = "The Autonomous Agent-to-Agent Micro-Economy"
p2.font.size = Pt(24)
p2.font.color.rgb = TEXT_LIGHT
p2.space_before = Pt(10)

p3 = tf1.add_paragraph()
p3.text = "International Hackathon Competition 2026 | Sofzenix IT Solution LLP"
p3.font.size = Pt(16)
p3.font.color.rgb = EMERALD_ACCENT
p3.space_before = Pt(30)

# SLIDE 2: The Problem
slide2 = add_blank_slide_with_bg(prs)
add_header(slide2, "The Real-World Problem: Friction in Agent Commerce")

problems = [
    ("Human-in-the-Loop Bottlenecks", "Modern AI agents executing multi-step tasks cannot independently procure resources—they stall waiting for human credit card approvals and API key setup."),
    ("The Micro-Payment Gateway Fee Trap", "Standard payment gateways charge fixed fees ($0.30 + 2.9%). A $0.001 agent call becomes mathematically unviable on traditional banking networks."),
    ("Lack of Trustless Settlement", "AI agents have no built-in wallet rails or escrow mechanisms to hold funds and release them only upon verified delivery of work.")
]

for i, (p_title, p_desc) in enumerate(problems):
    y_pos = Inches(1.8 + i * 1.7)
    card = slide2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), y_pos, Inches(11.73), Inches(1.4))
    card.fill.solid()
    card.fill.fore_color.rgb = CARD_BG
    card.line.color.rgb = MUTED_BORDER
    
    tf = card.text_frame
    tf.word_wrap = True
    p1 = tf.paragraphs[0]
    p1.text = f"• {p_title}"
    p1.font.size = Pt(18)
    p1.font.bold = True
    p1.font.color.rgb = CYAN_ACCENT
    
    p2 = tf.add_paragraph()
    p2.text = p_desc
    p2.font.size = Pt(14)
    p2.font.color.rgb = TEXT_LIGHT
    p2.space_before = Pt(5)

# SLIDE 3: The Solution
slide3 = add_blank_slide_with_bg(prs)
add_header(slide3, "The Solution: AgentX Sandbox Core Concept")

solutions = [
    ("Autonomous Agent Marketplace", "A simulated ecosystem where independent software agents hold wallets, negotiate task execution, and trade micro-services."),
    ("In-Memory Price Discovery", "Order-book matching engine where agents bid, ask, and dynamically price services based on latency, load, and SLA urgency."),
    ("Smart Contract Escrows", "Trustless escrow contracts lock native tokens ($CRED) and release payments instantly upon cryptographic or verification proof.")
]

for i, (s_title, s_desc) in enumerate(solutions):
    x_pos = Inches(0.8 + i * 3.9)
    card = slide3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x_pos, Inches(1.8), Inches(3.73), Inches(4.8))
    card.fill.solid()
    card.fill.fore_color.rgb = CARD_BG
    card.line.color.rgb = MUTED_BORDER
    
    tf = card.text_frame
    tf.word_wrap = True
    p1 = tf.paragraphs[0]
    p1.text = s_title
    p1.font.size = Pt(18)
    p1.font.bold = True
    p1.font.color.rgb = EMERALD_ACCENT
    
    p2 = tf.add_paragraph()
    p2.text = s_desc
    p2.font.size = Pt(14)
    p2.font.color.rgb = TEXT_LIGHT
    p2.space_before = Pt(15)

# SLIDE 4: Key Features & Technical Capabilities
slide4 = add_blank_slide_with_bg(prs)
add_header(slide4, "Key Features & Technical Capabilities")

rows, cols = 5, 3
table_shape = slide4.shapes.add_table(rows, cols, Inches(0.8), Inches(1.8), Inches(11.73), Inches(4.8))
table = table_shape.table

table.columns[0].width = Inches(2.8)
table.columns[1].width = Inches(3.5)
table.columns[2].width = Inches(5.43)

headers = ["Capability", "Engine Component", "Functional Impact"]
for i, h in enumerate(headers):
    cell = table.cell(0, i)
    cell.fill.solid()
    cell.fill.fore_color.rgb = CARD_BG
    p = cell.text_frame.paragraphs[0]
    p.text = h
    p.font.bold = True
    p.font.size = Pt(14)
    p.font.color.rgb = CYAN_ACCENT

data = [
    ("Agent Intelligence", "Groq Llama 3 70B + Rule Fallback", "Sub-second LLM decision making per tick with deterministic fallback rules."),
    ("Marketplace Engine", "In-Memory Price-Time Priority Matcher", "Matches buy/sell orders and micro-service requests with real-time price discovery."),
    ("Trustless Settlement", "TradeEscrow.sol Smart Contract", "Locks funds on-chain until task delivery is verified; refunds on timeout."),
    ("Live Observability", "React Flow Network Topology Graph", "Visualizes real-time agent node interactions, active escrows, and token transfers.")
]

for row_idx, row_data in enumerate(data, start=1):
    for col_idx, text in enumerate(row_data):
        cell = table.cell(row_idx, col_idx)
        cell.fill.solid()
        cell.fill.fore_color.rgb = DARK_BG
        p = cell.text_frame.paragraphs[0]
        p.text = text
        p.font.size = Pt(12)
        p.font.color.rgb = TEXT_LIGHT

# SLIDE 5: Technical Architecture
slide5 = add_blank_slide_with_bg(prs)
add_header(slide5, "Technical Architecture & Stack")

tech_blocks = [
    ("Frontend Layer", "React 18 • Vite • Tailwind CSS • Framer Motion • Lenis Scroll • Three.js 3D Canvas • React Flow Network Visualizer"),
    ("Backend & Orchestration", "Node.js • Express • Socket.io (2-Second Tick Loop) • In-Memory Order Book • Ethers.js Wallet Orchestrator"),
    ("AI & Reasoning Layer", "Groq API (Llama 3 70B Engine) • Low-Latency Decision Prompts • Deterministic Rule Strategy Fallback"),
    ("Blockchain & Settlement", "Polygon Amoy Testnet • AgentCredit.sol (ERC-20) • TradeEscrow.sol • ServiceRegistry.sol • Polygonscan Explorer")
]

for i, (b_title, b_desc) in enumerate(tech_blocks):
    y_pos = Inches(1.8 + i * 1.3)
    card = slide5.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), y_pos, Inches(11.73), Inches(1.1))
    card.fill.solid()
    card.fill.fore_color.rgb = CARD_BG
    card.line.color.rgb = MUTED_BORDER
    
    tf = card.text_frame
    tf.word_wrap = True
    p1 = tf.paragraphs[0]
    p1.text = b_title
    p1.font.size = Pt(16)
    p1.font.bold = True
    p1.font.color.rgb = CYAN_ACCENT
    
    p2 = tf.add_paragraph()
    p2.text = b_desc
    p2.font.size = Pt(13)
    p2.font.color.rgb = TEXT_LIGHT
    p2.space_before = Pt(4)

# SLIDE 6: Smart Contracts
slide6 = add_blank_slide_with_bg(prs)
add_header(slide6, "Smart Contracts & On-Chain Verification")

contracts = [
    ("AgentCredit.sol", "ERC-20 Token ($CRED)", "Serves as the native micro-currency minted to agent programmatic wallets upon simulation start."),
    ("TradeEscrow.sol", "Escrow Holding Vault", "Holds tokens in trust when a service order is placed. Releases payout upon verified completion."),
    ("ServiceRegistry.sol", "On-Chain Audit Log", "Logs service listings, execution receipts, and agent reliability metrics directly on Polygon Amoy.")
]

for i, (c_name, c_type, c_desc) in enumerate(contracts):
    x_pos = Inches(0.8 + i * 3.9)
    card = slide6.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x_pos, Inches(1.8), Inches(3.73), Inches(4.8))
    card.fill.solid()
    card.fill.fore_color.rgb = CARD_BG
    card.line.color.rgb = MUTED_BORDER
    
    tf = card.text_frame
    tf.word_wrap = True
    p1 = tf.paragraphs[0]
    p1.text = c_name
    p1.font.size = Pt(18)
    p1.font.bold = True
    p1.font.color.rgb = EMERALD_ACCENT
    
    p2 = tf.add_paragraph()
    p2.text = c_type
    p2.font.size = Pt(13)
    p2.font.bold = True
    p2.font.color.rgb = CYAN_ACCENT
    p2.space_before = Pt(5)
    
    p3 = tf.add_paragraph()
    p3.text = c_desc
    p3.font.size = Pt(13)
    p3.font.color.rgb = TEXT_LIGHT
    p3.space_before = Pt(15)

# SLIDE 7: UX & Interface Design
slide7 = add_blank_slide_with_bg(prs)
add_header(slide7, "Bespoke User Experience & Interface Design")

ux_highlights = [
    ("Bespoke 3D Landing Page", "Three.js interactive mesh background, smooth Lenis scroll transitions, and 'Get Started' CTA."),
    ("Interactive Topology Graph", "Real-time React Flow graph showing node connections, escrows, and token transfers."),
    ("Order Book & Price Chart", "Recharts spot price curve tracking alongside an active Bids/Asks depth table."),
    ("Transparent Agent Roster", "Agent cards displaying strategy profiles, balances, and live LLM reasoning logs.")
]

for i, (u_title, u_desc) in enumerate(ux_highlights):
    y_pos = Inches(1.6 + i * 1.3)
    card = slide7.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), y_pos, Inches(5.3), Inches(1.15))
    card.fill.solid()
    card.fill.fore_color.rgb = CARD_BG
    card.line.color.rgb = MUTED_BORDER
    
    tf = card.text_frame
    tf.word_wrap = True
    p1 = tf.paragraphs[0]
    p1.text = u_title
    p1.font.size = Pt(13)
    p1.font.bold = True
    p1.font.color.rgb = CYAN_ACCENT
    
    p2 = tf.add_paragraph()
    p2.text = u_desc
    p2.font.size = Pt(10.5)
    p2.font.color.rgb = TEXT_LIGHT
    p2.space_before = Pt(3)

# Add Generated User Flow Block Diagram Image on Right Column
image_path = r"C:\Users\sayim\.gemini\antigravity-ide\brain\9e16e32c-1886-4bd7-b45d-38b54cbf3c69\agentx_ui_user_flow_diagram_1786562643012.png"
import os
if os.path.exists(image_path):
    slide7.shapes.add_picture(image_path, Inches(6.3), Inches(1.6), Inches(6.2), Inches(5.2))

# SLIDE 8: Live Demo Controls
slide8 = add_blank_slide_with_bg(prs)
add_header(slide8, "Live Hackathon Demo & Control Suite")

demo_controls = [
    ("⚡ Seed Demo Scenario (3-Min)", "Triggers a deterministic 3-minute sequence showcasing market negotiation, micro-service escrow, and Polygonscan payout settlement."),
    ("💥 Inject Market Shock", "Forces a sudden surge in demand or service calls to demonstrate real-time dynamic pricing and order book depth adjustments."),
    ("📥 Export Run JSON", "Allows evaluators to download tick logs, state history, and smart contract transaction hashes for post-demo analysis.")
]

for i, (d_title, d_desc) in enumerate(demo_controls):
    y_pos = Inches(1.8 + i * 1.7)
    card = slide8.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), y_pos, Inches(11.73), Inches(1.4))
    card.fill.solid()
    card.fill.fore_color.rgb = CARD_BG
    card.line.color.rgb = MUTED_BORDER
    
    tf = card.text_frame
    tf.word_wrap = True
    p1 = tf.paragraphs[0]
    p1.text = d_title
    p1.font.size = Pt(18)
    p1.font.bold = True
    p1.font.color.rgb = EMERALD_ACCENT
    
    p2 = tf.add_paragraph()
    p2.text = d_desc
    p2.font.size = Pt(14)
    p2.font.color.rgb = TEXT_LIGHT
    p2.space_before = Pt(5)

# SLIDE 9: Business Impact & Roadmap
slide9 = add_blank_slide_with_bg(prs)
add_header(slide9, "Business Impact, Monetization & Roadmap")

impacts = [
    ("Real-World Use Cases", "• Idle GPU Compute Leasing\n• Automated Data Scraping & Oracles\n• Sub-task Delegation between AI Agents"),
    ("Protocol Monetization", "• Protocol Liquidity Fee (0.1% per trade)\n• Premium Registry Listings for Agents\n• Enterprise Sandbox Subscriptions"),
    ("Future Roadmap", "• Phase 1: Polygon Amoy Single Token\n• Phase 2: Cross-chain Multi-Asset Settlement\n• Phase 3: Plug-and-play LangChain SDK")
]

for i, (im_title, im_desc) in enumerate(impacts):
    x_pos = Inches(0.8 + i * 3.9)
    card = slide9.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x_pos, Inches(1.8), Inches(3.73), Inches(4.8))
    card.fill.solid()
    card.fill.fore_color.rgb = CARD_BG
    card.line.color.rgb = MUTED_BORDER
    
    tf = card.text_frame
    tf.word_wrap = True
    p1 = tf.paragraphs[0]
    p1.text = im_title
    p1.font.size = Pt(18)
    p1.font.bold = True
    p1.font.color.rgb = CYAN_ACCENT
    
    p2 = tf.add_paragraph()
    p2.text = im_desc
    p2.font.size = Pt(13)
    p2.font.color.rgb = TEXT_LIGHT
    p2.space_before = Pt(12)

# SLIDE 10: Conclusion & Links
slide10 = add_blank_slide_with_bg(prs)

title_box = slide10.shapes.add_textbox(Inches(1.0), Inches(1.5), Inches(11.3), Inches(4.5))
tf10 = title_box.text_frame
tf10.word_wrap = True

p1 = tf10.paragraphs[0]
p1.text = "AgentX Sandbox"
p1.font.size = Pt(44)
p1.font.bold = True
p1.font.color.rgb = CYAN_ACCENT

p2 = tf10.add_paragraph()
p2.text = "Empowering Autonomous Machine-to-Machine Commerce"
p2.font.size = Pt(20)
p2.font.color.rgb = TEXT_LIGHT
p2.space_before = Pt(10)

links = [
    ("Live Application:", "https://agentx-sandbox.vercel.app"),
    ("GitHub Repository:", "https://github.com/your-username/agentx-sandbox"),
    ("Network Settlement:", "Polygon Amoy Testnet (Chain ID: 80002)")
]

for label, val in links:
    p = tf10.add_paragraph()
    p.text = f"{label} {val}"
    p.font.size = Pt(14)
    p.font.color.rgb = EMERALD_ACCENT
    p.space_before = Pt(15)

# Save PowerPoint File
output_path = "AgentX_Sandbox_Hackathon_Deck_v2.pptx"
try:
    prs.save("AgentX_Sandbox_Hackathon_Deck.pptx")
    print("\nSUCCESS: Presentation generated at 'AgentX_Sandbox_Hackathon_Deck.pptx'")
except Exception:
    prs.save(output_path)
    print(f"\nSUCCESS: Presentation generated at '{output_path}'")